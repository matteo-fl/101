from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os


class StyleTemplate:
    """Класс для хранения стилей оформления"""

    def __init__(self, name, colors, fonts, layout_params):
        self.name = name
        self.colors = colors
        self.fonts = fonts
        self.layout_params = layout_params


def get_style_template(style_name: str) -> StyleTemplate:
    """Возвращает шаблон стиля по названию"""
    styles = {
        "corporate": StyleTemplate(
            name="Корпоративный",
            colors={
                "primary": RGBColor(0x00, 0x52, 0xCC),  # Синий
                "secondary": RGBColor(0x6C, 0x75, 0x7D),  # Серый
                "accent": RGBColor(0x00, 0x9C, 0xDF),  # Светло-синий
                "text": RGBColor(0x33, 0x33, 0x33),  # Тёмно-серый
                "background": RGBColor(0xFF, 0xFF, 0xFF)  # Белый
            },
            fonts={
                "title": Pt(36),
                "subtitle": Pt(24),
                "body": Pt(18),
                "small": Pt(14)
            },
            layout_params={
                "title_y": Inches(0.6),
                "content_y": Inches(1.8),
                "image_x": Inches(8.5),
                "image_y": Inches(1.5),
                "image_width": Inches(4)
            }
        ),
        "creative": StyleTemplate(
            name="Креативный",
            colors={
                "primary": RGBColor(0xFF, 0x6B, 0x6B),  # Коралловый
                "secondary": RGBColor(0x4E, 0xC9, 0xCA),  # Бирюзовый
                "accent": RGBColor(0xFF, 0xE6, 0x6D),  # Жёлтый
                "text": RGBColor(0x2C, 0x3E, 0x50),  # Тёмно-синий
                "background": RGBColor(0xF9, 0xF9, 0xF9)  # Светло-серый
            },
            fonts={
                "title": Pt(40),
                "subtitle": Pt(28),
                "body": Pt(20),
                "small": Pt(16)
            },
            layout_params={
                "title_y": Inches(0.8),
                "content_y": Inches(2.0),
                "image_x": Inches(8.0),
                "image_y": Inches(1.2),
                "image_width": Inches(4.5)
            }
        ),
        "minimalist": StyleTemplate(
            name="Минималистичный",
            colors={
                "primary": RGBColor(0x00, 0x00, 0x00),  # Чёрный
                "secondary": RGBColor(0x80, 0x80, 0x80),  # Серый
                "accent": RGBColor(0xC0, 0xC0, 0xC0),  # Светло-серый
                "text": RGBColor(0x33, 0x33, 0x33),  # Тёмно-серый
                "background": RGBColor(0xFF, 0xFF, 0xFF)  # Белый
            },
            fonts={
                "title": Pt(32),
                "subtitle": Pt(22),
                "body": Pt(16),
                "small": Pt(12)
            },
            layout_params={
                "title_y": Inches(0.5),
                "content_y": Inches(1.5),
                "image_x": Inches(9.0),
                "image_y": Inches(1.8),
                "image_width": Inches(3.5)
            }
        )
    }
    return styles.get(style_name, styles["corporate"])


def add_title_slide(prs, slide_data: dict, style: StyleTemplate):
    """Добавляет титульный слайд"""
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)

    # Настройка заголовка
    title = slide.shapes.title
    title.text = slide_data["title"]
    title.text_frame.paragraphs[0].font.size = style.fonts["title"]
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = style.colors["primary"]

    # Настройка подзаголовка
    if len(slide.placeholders) > 1:
        subtitle = slide.placeholders[1]
        subtitle.text = slide_data["content"]
        subtitle.text_frame.paragraphs[0].font.size = style.fonts["subtitle"]
        subtitle.text_frame.paragraphs[0].font.color.rgb = style.colors["secondary"]


def add_content_slide(prs, slide_data: dict, style: StyleTemplate, slide_num: int):
    """Добавляет контентный слайд"""
    layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(layout)

    # Заголовок слайда
    title_box = slide.shapes.add_textbox(
        Inches(0.8),
        style.layout_params["title_y"],
        Inches(11),
        Inches(1)
    )
    title_frame = title_box.text_frame
    title_frame.paragraphs[0].text = slide_data["title"]
    title_frame.paragraphs[0].font.size = style.fonts["title"]
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = style.colors["primary"]

    # Контент слайда
    content_box = slide.shapes.add_textbox(
        Inches(0.8),
        style.layout_params["content_y"],
        Inches(7),
        Inches(5)
    )
    content_frame = content_box.text_frame
    content_frame.word_wrap = True

    # Разбиваем контент на абзацы
    paragraphs = slide_data["content"].split('\n')
    for i, para_text in enumerate(paragraphs):
        if i == 0:
            p = content_frame.paragraphs[0]
        else:
            p = content_frame.add_paragraph()

        p.text = para_text
        p.font.size = style.fonts["body"]
        p.font.color.rgb = style.colors["text"]
        p.space_after = Pt(12)

    # Добавляем визуальные элементы в зависимости от стиля
    if style.name == "Креативный":
        # Добавляем декоративную линию
        line = slide.shapes.add_shape(
            1,  # Линия
            Inches(0.8),
            style.layout_params["title_y"] + Inches(0.8),
            Inches(3),
            Inches(0.05)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = style.colors["accent"]

    # Вставка изображения (если есть)
    img_path = f"uploads/img_{slide_num}.png"
    if os.path.exists(img_path):
        try:
            slide.shapes.add_picture(
                img_path,
                style.layout_params["image_x"],
                style.layout_params["image_y"],
                width=style.layout_params["image_width"]
            )
        except Exception as e:
            print(f"⚠️ Пропуск битой картинки для слайда {slide_num + 1}: {e}")


def add_layout_template(prs, slide_data: dict, style: StyleTemplate, slide_num: int, template_id: int):
    """Добавляет слайд с использованием специального шаблона разметки"""
    layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(layout)

    if template_id == 1:
        # Шаблон 1: Левая колонка с текстом, правая с изображением
        _add_template_two_columns(slide, slide_data, style, slide_num)
    elif template_id == 2:
        # Шаблон 2: Полноширинный текст с иконками/декоративными элементами
        _add_template_fullwidth(slide, slide_data, style, slide_num)
    elif template_id == 3:
        # Шаблон 3: Сетка (текст, изображение, дополнительное содержимое)
        _add_template_grid(slide, slide_data, style, slide_num)
    else:
        # Шаблон по умолчанию
        add_content_slide(prs, slide_data, style, slide_num)


def _add_template_two_columns(slide, slide_data: dict, style: StyleTemplate, slide_num: int):
    """Шаблон 1: Двухколончная разметка (текст слева, изображение справа)"""
    
    # Заголовок на всю ширину
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_frame.paragraphs[0].text = slide_data["title"]
    title_frame.paragraphs[0].font.size = style.fonts["title"]
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = style.colors["primary"]
    
    # Декоративная линия
    line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.4), Inches(2), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = style.colors["accent"]
    line.line.color.rgb = style.colors["accent"]
    
    # Левая колонка - текст
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.7), Inches(6.5), Inches(5.3))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    paragraphs = slide_data["content"].split('\n')
    for i, para_text in enumerate(paragraphs):
        if i == 0:
            p = content_frame.paragraphs[0]
        else:
            p = content_frame.add_paragraph()
        
        p.text = para_text.strip()
        p.font.size = style.fonts["body"]
        p.font.color.rgb = style.colors["text"]
        p.space_after = Pt(10)
        p.level = 0
    
    # Правая колонка - изображение
    img_path = f"uploads/img_{slide_num}.png"
    if os.path.exists(img_path):
        try:
            slide.shapes.add_picture(img_path, Inches(7.2), Inches(1.7), width=Inches(5.5))
        except Exception as e:
            print(f"⚠️ Ошибка вставки изображения: {e}")


def _add_template_fullwidth(slide, slide_data: dict, style: StyleTemplate, slide_num: int):
    """Шаблон 2: Полноширинный текст с декоративными элементами"""
    
    # Заголовок
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.733), Inches(1))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_frame.paragraphs[0].text = slide_data["title"]
    title_frame.paragraphs[0].font.size = style.fonts["title"]
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = style.colors["primary"]
    
    # Фоновый прямоугольник для выделения
    bg_shape = slide.shapes.add_shape(1, Inches(0.8), Inches(1.7), Inches(11.733), Inches(0.08))
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = style.colors["secondary"]
    bg_shape.line.color.rgb = style.colors["secondary"]
    
    # Основной текст
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(2), Inches(11.733), Inches(4.8))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    paragraphs = slide_data["content"].split('\n')
    for i, para_text in enumerate(paragraphs):
        if i == 0:
            p = content_frame.paragraphs[0]
        else:
            p = content_frame.add_paragraph()
        
        p.text = para_text.strip()
        p.font.size = style.fonts["body"]
        p.font.color.rgb = style.colors["text"]
        p.space_after = Pt(12)
        p.level = 0
    
    # Изображение внизу (если есть)
    img_path = f"uploads/img_{slide_num}.png"
    if os.path.exists(img_path):
        try:
            slide.shapes.add_picture(img_path, Inches(3), Inches(5), width=Inches(7.333))
        except Exception as e:
            print(f"⚠️ Ошибка вставки изображения: {e}")


def _add_template_grid(slide, slide_data: dict, style: StyleTemplate, slide_num: int):
    """Шаблон 3: Сетчатая разметка с блоками контента"""
    
    # Заголовок с фоном
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(1))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = style.colors["accent"]
    title_shape.line.color.rgb = style.colors["accent"]
    
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.25), Inches(11.733), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_frame.paragraphs[0].text = slide_data["title"]
    title_frame.paragraphs[0].font.size = style.fonts["title"]
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = style.colors["background"]
    
    # Верхний блок контента
    content_parts = slide_data["content"].split('\n')
    first_part = '\n'.join(content_parts[:len(content_parts)//2])
    second_part = '\n'.join(content_parts[len(content_parts)//2:])
    
    # Левый блок
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(6), Inches(5.8))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    
    for i, para_text in enumerate(first_part.split('\n')):
        if i == 0:
            p = left_frame.paragraphs[0]
        else:
            p = left_frame.add_paragraph()
        
        p.text = para_text.strip()
        p.font.size = style.fonts["body"]
        p.font.color.rgb = style.colors["text"]
        p.space_after = Pt(8)
    
    # Правый блок (если есть изображение, то здесь оно, иначе текст)
    img_path = f"uploads/img_{slide_num}.png"
    if os.path.exists(img_path):
        try:
            slide.shapes.add_picture(img_path, Inches(6.8), Inches(1.3), width=Inches(6))
        except Exception as e:
            print(f"⚠️ Ошибка вставки изображения: {e}")
    else:
        right_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.3), Inches(6), Inches(5.8))
        right_frame = right_box.text_frame
        right_frame.word_wrap = True
        
        for i, para_text in enumerate(second_part.split('\n')):
            if i == 0:
                p = right_frame.paragraphs[0]
            else:
                p = right_frame.add_paragraph()
            
            p.text = para_text.strip()
            p.font.size = style.fonts["body"]
            p.font.color.rgb = style.colors["text"]
            p.space_after = Pt(8)


def generate_presentation(slides_data: list, output_path: str, style_name: str = "corporate", template_id: int = 1) -> str:
    """
    Генерация презентации с выбранным стилем и шаблоном разметки

    Args:
        slides_data: список слайдов с полями title, content
        output_path: путь для сохранения
        style_name: стиль оформления (corporate, creative, minimalist)
        template_id: ID шаблона разметки (1, 2 или 3)
    """
    # Устанавливаем размеры слайда
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Получаем шаблон стиля
    style = get_style_template(style_name)

    print(f"🎨 Генерация презентации в стиле: {style.name}")
    print(f"📐 Используемый шаблон разметки: {template_id}")

    # Создаём слайды
    for i, slide in enumerate(slides_data):
        if i == 0:
            add_title_slide(prs, slide, style)
        else:
            if template_id in [1, 2, 3]:
                add_layout_template(prs, slide, style, i, template_id)
            else:
                add_content_slide(prs, slide, style, i)

    # Сохраняем презентацию
    prs.save(output_path)
    print(f"✅ Презентация сохранена: {output_path}")

    return output_path